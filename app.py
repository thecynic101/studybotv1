import streamlit as st
import google.generativeai as genai
import PyPDF2
from PIL import Image
import docx # For Word Docs
from pptx import Presentation # For PowerPoints
import io

# --- CONFIGURATION ---
st.set_page_config(page_title="StudyBot Pro", page_icon="🎓", layout="wide")

# --- AUTHENTICATION ---
api_key = None
if "GEMINI_API_KEY" in st.secrets:
    api_key = st.secrets["GEMINI_API_KEY"]

if not api_key:
    api_key = st.sidebar.text_input("⚠️ Admin Mode: Enter API Key", type="password")

# --- LOGIC FUNCTIONS ---
def extract_text_from_files(uploaded_files):
    full_text = ""
    
    for file in uploaded_files:
        file_name = file.name.lower()
        
        # 1. Handle PDF
        if file_name.endswith('.pdf'):
            try:
                pdf_reader = PyPDF2.PdfReader(file)
                for i, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text:
                        full_text += f"\n\n--- Source: {file.name}, Page {i+1} ---\n{text}"
            except Exception as e:
                st.error(f"Error reading PDF {file.name}: {e}")

        # 2. Handle Word (.docx)
        elif file_name.endswith('.docx'):
            try:
                doc = docx.Document(file)
                text = "\n".join([para.text for para in doc.paragraphs])
                full_text += f"\n\n--- Source: {file.name} (Full Doc) ---\n{text}"
            except Exception as e:
                st.error(f"Error reading Word Doc {file.name}: {e}")

        # 3. Handle PowerPoint (.pptx)
        elif file_name.endswith('.pptx'):
            try:
                prs = Presentation(file)
                for i, slide in enumerate(prs.slides):
                    slide_text = ""
                    # Extract text from shapes in the slide
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text += shape.text + " "
                    full_text += f"\n\n--- Source: {file.name}, Slide {i+1} ---\n{slide_text}"
            except Exception as e:
                st.error(f"Error reading PowerPoint {file.name}: {e}")
                
    return full_text

def get_gemini_response(api_key, study_material, image_input):
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel('gemini-1.5-flash')

    prompt = """
    You are an expert academic tutor.
    
    TASK:
    1. Analyze the exam questions in the image.
    2. Answer them using ONLY the provided Study Material.
    3. For every answer, you MUST cite the (Source, Page/Slide Number).
    4. If the answer is not in the material, state "Not found in provided notes."
    
    FORMAT:
    Question 1: [Question Text]
    Answer: [Detailed Answer]
    📍 Citation: [Source, Page/Slide X]
    
    Study Material:
    """
    
    try:
        response = model.generate_content([prompt + study_material, image_input])
        return response.text
    except Exception as e:
        if "429" in str(e):
            return "⚠️ SERVER BUSY: Too many students are studying right now. Please wait 1 minute and try again."
        else:
            return f"Error: {e}"

# --- UI ---
st.title("🎓 StudyBot: Automate Your Revision")
st.markdown("Upload Handouts (PDF, PPTX, DOCX) + Past Questions -> Get a Cited Answer Sheet.")

if 'answer_sheet' not in st.session_state:
    st.session_state['answer_sheet'] = None

col1, col2 = st.columns([1, 1])

with col1:
    st.subheader("1. Your Materials")
    # Updated to accept all 3 formats
    uploaded_files = st.file_uploader(
        "Upload Handouts", 
        type=['pdf', 'docx', 'pptx'], 
        accept_multiple_files=True
    )
    
    st.subheader("2. The Exam")
    question_image = st.file_uploader("Upload Question Paper (Image)", type=['png', 'jpg', 'jpeg'])
    
    generate_btn = st.button("🚀 Generate Cheat Sheet", type="primary", use_container_width=True)

with col2:
    st.subheader("3. The Answers")
    
    if generate_btn:
        if not api_key:
            st.error("API Key missing. Contact Admin.")
        elif not uploaded_files or not question_image:
            st.warning("Please upload both handouts and questions.")
        else:
            with st.spinner("Reading files... (This might take a moment)"):
                # 1. Process Files (PDF/DOCX/PPTX)
                raw_text = extract_text_from_files(uploaded_files)
                
                if not raw_text:
                    st.error("Could not extract text. Please check your files.")
                else:
                    # 2. Process Image
                    img = Image.open(question_image)
                    # 3. AI Magic
                    result = get_gemini_response(api_key, raw_text, img)
                    
                    st.session_state['answer_sheet'] = result
                    st.rerun()

    if st.session_state['answer_sheet']:
        st.markdown(st.session_state['answer_sheet'])
        
        st.download_button(
            label="📥 Download Answer Sheet (.txt)",
            data=st.session_state['answer_sheet'],
            file_name="studybot_answers.txt",
            mime="text/plain",
            use_container_width=True
        )
